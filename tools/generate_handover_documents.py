"""Generate the formal report and presentation (optional doc tooling, not app runtime)."""
from pathlib import Path
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as C
from pptx.util import Inches as I, Pt as P

ROOT = Path(__file__).resolve().parents[1]
REPO = "https://github.com/passantabdou23-svg/Construction-MultiAgent-System"
NAVY, TEAL, ORANGE, WHITE = C(24,39,75), C(0,137,123), C(238,124,43), C(255,255,255)

SECTIONS = [
 ("Executive summary", [
  "This project implements a governed construction-change workflow: a site note is validated, supported by controlled regulatory evidence, converted to an immutable package, and held for an authenticated human decision. Procurement planning and Critical Path Method (CPM) analysis remain blocked until approval.",
  "The accepted build combines local Ollama inference with deterministic validation, persistent Retrieval-Augmented Generation (RAG), role-based access control, separation of duties, SQLite transactions, and a hash-chained audit trail. Its purpose is controlled planning support—not autonomous engineering approval." ]),
 ("1. Problem, objective, and scope", [
  "Ungoverned agent output can create unsupported technical claims, premature purchasing, and untraceable decisions. The objective is evidence-backed multi-agent coordination with explicit human authorization before consequential actions.",
  "The operating scope is a controlled single-workstation prototype. Licensed engineers and applicable authorities remain responsible for engineering and compliance decisions." ]),
 ("2. Implemented architecture", [
  "Interface: authenticated, role-aware Streamlit control centre and read-only audit views.",
  "Orchestration: sequential agent pipeline with explicit run states and approval gate.",
  "Knowledge: four controlled England Approved Documents, clause-aware chunks, pretrained MiniLM embeddings, and persistent ChromaDB.",
  "Safety: Pydantic contracts, confidence rejection, exact-quote and numeric support checks, and document-version governance.",
  "Persistence: transactional SQLite records, ordered SHA-256 audit chain, and validated backups." ]),
 ("3. Governed workflow", [
  "A PREPARER submits a bounded site note. Validation rejects unsafe inputs before model invocation. Discipline-aware hybrid RAG retrieves controlled evidence. The design agent proposes structured claims; each accepted claim must cite an immutable chunk and verbatim quote.",
  "An immutable snapshot enters the queue. A different authorized reviewer reauthenticates and decides exactly once. Only approval permits procurement estimation and deterministic CPM analysis. Every state change is persisted and hash-chained." ]),
 ("4. Controlled RAG implementation", [
  "Corpus: official England Approved Documents A, C, K, and 7; 198 PDF pages, 190 searchable pages; 586 auditable chunks and 526 retrieval-eligible chunks.",
  "Retrieval: pretrained sentence-transformers/all-MiniLM-L6-v2 running locally; persistent ChromaDB; 75% semantic similarity plus 25% BM25-style lexical relevance; calibrated confidence floor 0.45.",
  "Citations preserve document, section, clause, printed page, PDF page, edition, status, jurisdiction, checksum, and official source URL. Covers and contents remain auditable but are excluded from retrieval." ]),
 ("5. Grounding and deterministic safeguards", [
  "The guard rejects unsupported numbers, unavailable chunks, altered quotes, incompatible versions, and evidence-created procurement items. Site facts are isolated from regulatory evidence. A bounded correction attempt may repair a malformed model contract, but deterministic verification remains authoritative.",
  "Two live Ollama trials introduced unsupported numeric claims; both packages were refused before approval, procurement, or scheduling." ]),
 ("6. Authentication, authorization, and audit", [
  "Passwords use unique salts and scrypt derivation. Server-side roles are PREPARER, DESIGN_REVIEWER, PROJECT_MANAGER, and ADMIN. Five failed logins trigger temporary lockout; sessions have idle and absolute expiry.",
  "Separation of duties blocks self-approval. Decisions require fresh password verification. Immutable snapshots and single-use decisions block stale-state approval and replay. The ordered SHA-256 chain detects ordinary modification, deletion, insertion, and reordering." ]),
 ("7. Procurement and CPM boundaries", [
  "Procurement outputs are unverified planning estimates. No live supplier API, verified quotation, or automatic purchase is claimed. Python recalculates totals and dates instead of trusting model arithmetic.",
  "NetworkX computes CPM effects over a five-task demonstration network. This is not a live Primavera P6 or Microsoft Project integration." ]),
 ("8. Verification results", [
  "73 repository tests and 3 stress tests passed. Retrieval Hit@5 and Top-1 accuracy were 100%; Mean Reciprocal Rank was 1.000; document routing, positive acceptance, negative rejection, grounded supported-case acceptance, and grounding-guard rejection were 100%.",
  "These percentages describe only the labelled controlled evaluation sets. They are not a universal construction-compliance benchmark.",
  "Accepted workflow evidence: 3 active accounts, 1 grounded revision, 1 approved immutable package, 1 unverified procurement record, 1 CPM impact, 0 pending approvals, 18 valid audit events, and a validated pre-release backup." ]),
 ("9. Deployment and recovery", [
  "Use Python 3.12 and exact dependency pins. Run ingestion, indexing, retrieval evaluation, grounding evaluation, repository tests, stress tests, and release_check.py. Generate ACCEPTANCE_REPORT.md and ACCEPTANCE_EVIDENCE.json against a validated backup.",
  "Bind Streamlit to localhost by default. LAN credentials require HTTPS. Backups use SQLite online backup and SHA-256 manifests; guarded restore verifies a disposable copy, preserves the old destination, and replaces atomically." ]),
 ("10. Limitations and future work", [
  "No multi-factor authentication (MFA), enterprise single sign-on (SSO), account recovery, legally binding electronic signature, verified supplier integration, or production monitoring service. The audit chain is tamper-evident but not digitally signed or externally anchored.",
  "Approved Documents apply to England and are not Egyptian compliance authority. Future work includes external identity, signed audit anchoring, live verified connectors, larger independent evaluations, managed observability, and formal licensed-engineer studies." ]),
 ("11. Conclusion", [
  "The accepted release is a coherent, reproducible, and defensively bounded construction multi-agent prototype. Its main contribution is controlled evidence plus deterministic refusal, authenticated human approval, and traceable state transitions.",
  f"Repository: {REPO}" ]),
]

SLIDES = [
 ("Problem and objective", ["Construction changes create technical, procurement, and schedule consequences.","Ungoverned LLM output can create unsupported claims or premature actions.","Objective: controlled evidence, deterministic checks, human authority, and traceability.","Success means safe state transitions—not autonomous approval."], ORANGE),
 ("Scope and trust boundary", ["Controlled single-workstation prototype using local Streamlit, Ollama, ChromaDB, and SQLite.","Approved Documents A, C, K, and 7 apply to England—not Egyptian compliance authority.","Procurement outputs are unverified planning estimates.","CPM uses a five-task demonstration network.","Licensed engineers retain decision responsibility."], ORANGE),
 ("Implemented architecture", ["Streamlit → validation → governed agent pipeline","Discipline router → hybrid retrieval → controlled citations","Ollama extraction → deterministic grounding guard","Immutable package → authenticated human approval","Approval → procurement planning + deterministic CPM","SQLite → SHA-256 audit chain → validated backup"], TEAL),
 ("Governed workflow", ["1  PREPARER submits a bounded site note","2  RAG retrieves controlled versioned evidence","3  Grounding verifies claims, numbers, chunks, and exact quotes","4  Immutable snapshot enters the approval queue","5  A different reviewer reauthenticates and decides once","6  Approval alone unlocks procurement and CPM","7  Every transition is persisted and hash-chained"], TEAL),
 ("Controlled RAG corpus", ["4 official England Approved Documents: A, C, K, and 7","198 PDF pages; 190 searchable pages","586 auditable chunks; 526 retrieval-eligible chunks","Clause-level citation and governance metadata","Covers/contents audited but excluded from retrieval","Source and embedding contracts prevent stale reuse"], TEAL),
 ("Persistent hybrid retrieval", ["Pretrained all-MiniLM-L6-v2 embeddings run locally","Persistent ChromaDB survives restarts","75% semantic + 25% BM25-style lexical relevance","Discipline-aware routing reduces cross-document noise","Confidence floor 0.45","Low-confidence and out-of-scope questions are refused"], TEAL),
 ("Claim-level grounding", ["Every accepted claim cites an immutable chunk and verbatim quote.","Numbers must be supported by evidence or bounded site facts.","Unavailable chunks, altered quotes, conflicts, and evidence-created items are rejected.","Deterministic verification remains authoritative.","Unsupported numeric claims in two live trials were refused safely."], ORANGE),
 ("Authenticated human approval", ["PREPARER · DESIGN_REVIEWER · PROJECT_MANAGER · ADMIN","Scrypt passwords, server-side sessions, lockout and expiry","Separation of duties blocks self-approval","Fresh password verification before decision","Immutable snapshots and single-use decisions block replay","Procurement and scheduling remain blocked until approval"], ORANGE),
 ("Procurement and CPM boundaries", ["Supplier, price, and delivery records remain PENDING_VERIFICATION.","No live supplier API or verified quotation source is claimed.","Python recalculates totals and dates.","NetworkX computes deterministic CPM after approval.","The schedule is a five-task demonstration—not a production plan."], ORANGE),
 ("Data integrity and audit", ["Transactional SQLite and non-destructive migrations","Ordered SHA-256 event chain","Server-side integrity verification","SQLite online backup + SHA-256 manifest","Guarded, validated, atomic restore"], TEAL),
 ("Accepted demonstration", ["3 active authenticated accounts","1 grounded revision and immutable package","Different reviewer approved after reauthentication","1 unverified procurement record and 1 CPM impact","0 pending approvals and 0 failed runs","18-event audit chain and pre-release backup validated"], TEAL),
 ("Measured verification", ["73 repository tests + 3 stress tests passed","Hit@5 100% · Top-1 100% · MRR 1.000","Routing 100%","Positive acceptance 100% · negative rejection 100%","Grounding acceptance 100% · guard rejection 100%","Metrics apply only to labelled controlled sets"], TEAL),
 ("Deployment and recovery", ["Python 3.12 with exact dependency pins","Localhost default on the Ollama workstation","LAN credentials require HTTPS","Deterministic release and acceptance checks","Privacy-safe Markdown + JSON evidence","Activation-ready GitHub Actions template"], TEAL),
 ("Explicit limitations", ["No MFA, enterprise SSO, recovery, or legal e-signature","Audit is not digitally signed or externally anchored","No verified supplier integration or production monitoring","England sources do not establish Egyptian compliance","Limited evaluation sets do not prove universal correctness","System supports—not replaces—licensed review"], ORANGE),
 ("Future work", ["External identity, MFA, recovery, and enterprise RBAC","Signed or externally anchored audit checkpoints","Verified supplier and project-planning connectors","Larger independent and adversarial evaluations","Production observability and incident response","Licensed-engineer safety and usability studies"], ORANGE),
 ("Project assets and demonstration", ["Live workflow: prepare → inspect → reauthenticate → decide → trace","Role-aware dashboard and audit views","ACCEPTANCE_REPORT.md + ACCEPTANCE_EVIDENCE.json","mULTIaGENTS.docx","MultiAgents.mp4",f"Repository: {REPO}"], TEAL),
 ("Conclusion", ["A complete governed prototype—not an autonomous construction authority.","Controlled evidence + deterministic refusal + authenticated approval + traceability.","Reproducible, testable, recoverable, and explicit about limitations.","Questions and review feedback are welcome."], NAVY),
]

def report():
 d=Document(); s=d.sections[0]; s.top_margin=s.bottom_margin=Inches(.65); s.left_margin=s.right_margin=Inches(.8)
 d.styles["Normal"].font.name="Aptos"; d.styles["Normal"].font.size=Pt(10.5)
 p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run("Construction Multi-Agent Control Centre"); r.bold=True; r.font.size=Pt(25); r.font.color.rgb=RGBColor(24,39,75)
 p=d.add_paragraph("Final Technical Report and Acceptance Handover"); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
 p=d.add_paragraph(f"Python 3.12 · Streamlit · Ollama · Hybrid RAG · SQLite\n{REPO}"); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
 d.add_page_break()
 for title,paras in SECTIONS:
  h=d.add_heading(title,level=1); h.style.font.color.rgb=RGBColor(24,39,75)
  for text in paras: d.add_paragraph(text)
 d.save(ROOT/"mULTIaGENTS.docx")

def slide(prs,title,bullets,accent):
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
 bar=s.shapes.add_shape(1,0,0,I(13.333),I(.16)); bar.fill.solid(); bar.fill.fore_color.rgb=accent; bar.line.fill.background()
 t=s.shapes.add_textbox(I(.65),I(.42),I(12),I(.7)); p=t.text_frame.paragraphs[0]; p.text=title; p.font.size=P(28); p.font.bold=True; p.font.color.rgb=NAVY
 b=s.shapes.add_textbox(I(.85),I(1.4),I(11.7),I(5.4)); tf=b.text_frame; tf.clear(); tf.word_wrap=True
 for n,x in enumerate(bullets):
  p=tf.paragraphs[0] if n==0 else tf.add_paragraph(); p.text=x; p.font.size=P(18); p.font.color.rgb=NAVY; p.space_after=P(12)
 f=s.shapes.add_textbox(I(.7),I(7.05),I(12),I(.2)); p=f.text_frame.paragraphs[0]; p.text="Controlled local prototype · Acronyms are defined at first use or in context"; p.font.size=P(8); p.font.color.rgb=C(100,110,120)

def presentation():
 prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
 s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
 t=s.shapes.add_textbox(I(.85),I(1.3),I(11.7),I(2)); p=t.text_frame.paragraphs[0]; p.text="Construction Multi-Agent\nControl Centre"; p.font.size=P(38); p.font.bold=True; p.font.color.rgb=WHITE
 u=s.shapes.add_textbox(I(.9),I(3.85),I(10.9),I(1.6)); u.text_frame.word_wrap=True; p=u.text_frame.paragraphs[0]; p.text="Governed local AI with controlled RAG, deterministic grounding, authenticated approval, CPM analysis, and auditable persistence"; p.font.size=P(17); p.font.color.rgb=C(205,225,238)
 r=s.shapes.add_textbox(I(.9),I(6.55),I(11.5),I(.35)); p=r.text_frame.paragraphs[0]; p.text=REPO; p.font.size=P(10); p.font.color.rgb=WHITE
 for title,bullets,accent in SLIDES: slide(prs,title,bullets,accent)
 prs.save(ROOT/"Construction Multi-Agent System Presentation.pptx")

if __name__=="__main__": report(); presentation(); print("Generated formal Word report and PowerPoint presentation.")
